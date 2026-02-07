"""Document processors for DOCX, PPTX, TXT, XLSX files"""
import logging
from typing import List, Dict, Any
from pathlib import Path
import docx
from pptx import Presentation
import pandas as pd
from openpyxl import load_workbook

logger = logging.getLogger(__name__)


async def extract_text_from_docx(file_path: str) -> List[Dict[str, Any]]:
    """
    Extract text from DOCX file
    
    Args:
        file_path: Path to DOCX file
        
    Returns:
        List of dicts with {text, metadata}
    """
    try:
        doc = docx.Document(file_path)
        paragraphs = []
        
        logger.info(f"Processing DOCX file: {file_path}")
        
        for idx, paragraph in enumerate(doc.paragraphs):
            text = paragraph.text.strip()
            if text:
                paragraphs.append({
                    'text': text,
                    'metadata': {
                        'paragraph_index': idx,
                        'source_type': 'docx'
                    }
                })
        
        # Also extract text from tables
        for table_idx, table in enumerate(doc.tables):
            table_text = []
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_text.append(cell_text)
                if row_text:
                    table_text.append(' | '.join(row_text))
            
            if table_text:
                paragraphs.append({
                    'text': '\n'.join(table_text),
                    'metadata': {
                        'table_index': table_idx,
                        'source_type': 'docx_table'
                    }
                })
        
        logger.info(f"Extracted {len(paragraphs)} text segments from DOCX")
        return paragraphs
        
    except Exception as e:
        logger.error(f"Error processing DOCX file {file_path}: {str(e)}")
        raise Exception(f"Failed to process DOCX: {str(e)}")


async def extract_text_from_pptx(file_path: str) -> List[Dict[str, Any]]:
    """
    Extract text from PPTX file with slide numbers
    
    Args:
        file_path: Path to PPTX file
        
    Returns:
        List of dicts with {text, slide_number, metadata}
    """
    try:
        prs = Presentation(file_path)
        slides_data = []
        
        logger.info(f"Processing PPTX with {len(prs.slides)} slides: {file_path}")
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = []
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                slides_data.append({
                    'text': '\n'.join(slide_text),
                    'slide_number': slide_num,
                    'metadata': {
                        'total_slides': len(prs.slides),
                        'slide_number': slide_num,
                        'source_type': 'pptx'
                    }
                })
        
        logger.info(f"Extracted text from {len(slides_data)} slides")
        return slides_data
        
    except Exception as e:
        logger.error(f"Error processing PPTX file {file_path}: {str(e)}")
        raise Exception(f"Failed to process PPTX: {str(e)}")


async def extract_text_from_txt(file_path: str) -> List[Dict[str, Any]]:
    """
    Extract text from plain text file
    
    Args:
        file_path: Path to TXT file
        
    Returns:
        List with single dict containing {text, metadata}
    """
    try:
        logger.info(f"Processing TXT file: {file_path}")
        
        with open(file_path, 'r', encoding='utf-8') as file:
            text = file.read()
        
        if not text.strip():
            logger.warning(f"Empty text file: {file_path}")
            return []
        
        return [{
            'text': text.strip(),
            'metadata': {
                'source_type': 'txt',
                'file_size': len(text)
            }
        }]
        
    except UnicodeDecodeError:
        # Try with different encoding
        try:
            with open(file_path, 'r', encoding='latin-1') as file:
                text = file.read()
            
            return [{
                'text': text.strip(),
                'metadata': {
                    'source_type': 'txt',
                    'encoding': 'latin-1',
                    'file_size': len(text)
                }
            }]
        except Exception as e:
            logger.error(f"Error processing TXT file with alternate encoding: {str(e)}")
            raise Exception(f"Failed to process TXT: {str(e)}")
    except Exception as e:
        logger.error(f"Error processing TXT file {file_path}: {str(e)}")
        raise Exception(f"Failed to process TXT: {str(e)}")


async def extract_text_from_xlsx(file_path: str) -> List[Dict[str, Any]]:
    """
    Extract text from XLSX file
    
    Args:
        file_path: Path to XLSX file
        
    Returns:
        List of dicts with {text, sheet_name, metadata}
    """
    try:
        logger.info(f"Processing XLSX file: {file_path}")
        
        # Load workbook
        wb = load_workbook(file_path, read_only=True)
        sheets_data = []
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            
            # Read sheet data
            sheet_text = []
            for row in ws.iter_rows(values_only=True):
                row_text = [str(cell) for cell in row if cell is not None and str(cell).strip()]
                if row_text:
                    sheet_text.append(' | '.join(row_text))
            
            if sheet_text:
                sheets_data.append({
                    'text': '\n'.join(sheet_text),
                    'sheet_name': sheet_name,
                    'metadata': {
                        'total_sheets': len(wb.sheetnames),
                        'sheet_name': sheet_name,
                        'source_type': 'xlsx'
                    }
                })
        
        wb.close()
        logger.info(f"Extracted text from {len(sheets_data)} sheets")
        return sheets_data
        
    except Exception as e:
        logger.error(f"Error processing XLSX file {file_path}: {str(e)}")
        raise Exception(f"Failed to process XLSX: {str(e)}")
